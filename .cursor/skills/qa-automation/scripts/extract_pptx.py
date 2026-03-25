#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPTX 기획서에서 텍스트·표(markitdown), 노트·이미지(python-pptx)를 추출하여 공통 스키마로 반환.
Usage: python extract_pptx.py <pptx_path> [--reference-dir DIR]
"""

import hashlib
import os
import re
import sys
from concurrent.futures import ThreadPoolExecutor

from _utils import parse_reference_dir, resolve_ref_dir, setup_stdout_utf8, to_rel_path

CONTENT_TYPE_TO_EXT = {"image/jpeg": "jpg", "image/png": "png", "image/gif": "gif", "image/bmp": "bmp"}


def _parse_markitdown_slides(md_text):
    """markitdown PPTX 출력을 슬라이드 번호 → markdown 내용으로 분리."""
    slides = {}
    current_slide = None
    current_lines = []

    for line in md_text.splitlines():
        m = re.match(r"<!--\s*Slide number:\s*(\d+)\s*-->", line)
        if m:
            if current_slide is not None:
                slides[current_slide] = "\n".join(current_lines).strip()
            current_slide = int(m.group(1))
            current_lines = []
        elif current_slide is not None:
            current_lines.append(line)

    if current_slide is not None:
        slides[current_slide] = "\n".join(current_lines).strip()

    return slides


def _clean_markdown(text):
    """마크다운 마커를 제거하여 순수 텍스트로 변환."""
    text = re.sub(r"^#+\s*", "", text).strip()
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)  # [text](url) → text
    text = re.sub(r"^[-*+]\s+", "", text)  # 리스트 마커
    text = re.sub(r"^\d+\.\s+", "", text)  # 번호 리스트
    return text.strip()


def _extract_texts_and_tables(md_content):
    """슬라이드 markdown에서 텍스트 목록과 표(2차원 배열) 목록을 추출."""
    lines = md_content.splitlines()
    texts = []
    tables = []
    i = 0

    while i < len(lines):
        line = lines[i]
        if line.strip().startswith("|"):
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i])
                i += 1
            rows = []
            for tl in table_lines:
                if re.match(r"\s*\|[-|\s:]+\|\s*$", tl):
                    continue
                cells = [c.strip() for c in tl.strip().strip("|").split("|")]
                rows.append(cells)
            if rows:
                tables.append(rows)
        else:
            text = _clean_markdown(line)
            if text and text != "---":
                texts.append(text)
            i += 1

    return texts, tables


def _get_shape_alt_text(shape):
    """셰이프에서 alt text(대체 텍스트)를 추출."""
    try:
        desc = shape._element.attrib.get("descr", "")
        if desc:
            return desc
        nvSpPr = shape._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetml}cNvPr")
        if nvSpPr is None:
            for child in shape._element.iter():
                if child.tag.endswith("}cNvPr"):
                    desc = child.attrib.get("descr", "")
                    if desc:
                        return desc
    except Exception:
        pass
    return ""


def _collect_images_recursive(shape, MSO_SHAPE_TYPE):
    """셰이프에서 이미지를 재귀적으로 수집 (그룹 셰이프 포함)."""
    images = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            images.extend(_collect_images_recursive(child, MSO_SHAPE_TYPE))
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        images.append(shape)
    return images


def _extract_texts_from_shapes(slide):
    """python-pptx로 슬라이드의 텍스트를 직접 추출 (markitdown fallback)."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
    return texts


def _run_markitdown(pptx_path):
    """Thread 1: markitdown으로 텍스트·표 추출."""
    try:
        from markitdown import MarkItDown
        md_result = MarkItDown().convert(pptx_path)
        return _parse_markitdown_slides(md_result.text_content)
    except Exception as e:
        print(f"[WARN] markitdown 변환 실패, python-pptx fallback 사용: {e}")
        return {}


def _run_pptx_extraction(pptx_path, ref_dir):
    """Thread 2: python-pptx로 이미지·노트·텍스트 추출."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx_path)
    slides_data = []
    seen_hashes = {}

    for slide_idx, slide in enumerate(prs.slides, 1):
        texts = _extract_texts_from_shapes(slide)

        images = []
        ref_images = []
        for shape in slide.shapes:
            for pic_shape in _collect_images_recursive(shape, MSO_SHAPE_TYPE):
                try:
                    img = pic_shape.image
                    blob = img.blob
                    content_hash = hashlib.md5(blob).hexdigest()

                    if content_hash in seen_hashes:
                        rel_path = seen_hashes[content_hash]
                        alt_text = _get_shape_alt_text(pic_shape)
                        images.append({"path": rel_path, "description": alt_text})
                        ref_images.append({"source_page": slide_idx, "path": rel_path})
                        continue

                    ext = CONTENT_TYPE_TO_EXT.get(img.content_type, "png") or "png"
                    rel_name = f"slide_{slide_idx}_img_{len(images)}.{ext}"
                    abs_path = os.path.join(ref_dir, rel_name)
                    with open(abs_path, "wb") as f:
                        f.write(blob)
                    rel_path = to_rel_path(abs_path, os.path.join("outputs", "reference", rel_name))
                    seen_hashes[content_hash] = rel_path
                    alt_text = _get_shape_alt_text(pic_shape)
                    images.append({"path": rel_path, "description": alt_text})
                    ref_images.append({"source_page": slide_idx, "path": rel_path})
                except Exception as e:
                    images.append({"path": "", "description": f"(추출 실패: {e})"})

        notes = ""
        if slide.has_notes_slide:
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text:
                    notes = notes_frame.text.strip()
            except Exception:
                pass

        slides_data.append({
            "slide_idx": slide_idx,
            "texts": texts,
            "images": images,
            "ref_images": ref_images,
            "notes": notes,
        })

    return slides_data


def extract_pptx(pptx_path, reference_dir=None):
    """PPTX 파일에서 내용 추출. 공통 스키마 반환."""
    try:
        from markitdown import MarkItDown  # noqa: F401
    except ImportError:
        print("Error: markitdown 라이브러리가 필요합니다.")
        print("설치: pip install markitdown")
        return None
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        print("Error: python-pptx 라이브러리가 필요합니다.")
        print("설치: pip install python-pptx")
        return None

    if not os.path.exists(pptx_path):
        print(f"Error: 파일을 찾을 수 없습니다 - {pptx_path}")
        return None

    ref_dir = resolve_ref_dir(reference_dir)

    # markitdown과 python-pptx를 병렬 실행
    with ThreadPoolExecutor(max_workers=2) as executor:
        md_future = executor.submit(_run_markitdown, pptx_path)
        pptx_future = executor.submit(_run_pptx_extraction, pptx_path, ref_dir)

        slide_md_map = md_future.result()
        pptx_data = pptx_future.result()

    # 병합: markitdown 텍스트/표가 있으면 우선, 없으면 python-pptx 텍스트 사용
    slides_out = []
    reference_images = []

    for data in pptx_data:
        slide_idx = data["slide_idx"]
        md_content = slide_md_map.get(slide_idx, "")
        texts, tables = _extract_texts_and_tables(md_content)
        if not texts:
            texts = data["texts"]

        slides_out.append({
            "page_num": slide_idx,
            "texts": texts,
            "tables": tables,
            "notes": data["notes"],
            "images": data["images"],
        })
        reference_images.extend(data["ref_images"])

    result = {"slides": slides_out, "reference_images": reference_images}

    print(f"=== {os.path.basename(pptx_path)} 분석 결과 ===\n")
    for slide in slides_out:
        print(f"--- Slide {slide['page_num']} ---")
        for text in slide["texts"]:
            print(text)
        if slide["tables"]:
            print("\n[표 데이터]")
            for table in slide["tables"]:
                for row in table:
                    print(" | ".join(row))
                print()
        if slide["notes"]:
            print(f"\n[노트] {slide['notes']}")
        if slide["images"]:
            print("\n[삽입 이미지]")
            for im in slide["images"]:
                if im["path"]:
                    desc = f' ({im["description"]})' if im["description"] else ""
                    print(f"  - {im['path']}{desc}")
        print()

    return result


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <pptx_path> [--reference-dir DIR]")
        print("Example: python extract_pptx.py inputs/quiz.pptx")
        sys.exit(1)

    setup_stdout_utf8()
    r = extract_pptx(sys.argv[1], reference_dir=parse_reference_dir(sys.argv))
    sys.exit(0 if r is not None else 1)


if __name__ == "__main__":
    main()
