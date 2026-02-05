import os
from pptx import Presentation

# 공통 AI 설정 모듈에서 model import
from ai_config import model

class DocAnalyst:
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.raw_content = []

    def extract_text_and_tables(self):
        """PPTX 슬라이드에서 모든 텍스트와 표 데이터를 추출"""
        prs = Presentation(self.pptx_path)
        
        for i, slide in enumerate(prs.slides):
            slide_data = {"slide_index": i + 1, "text": [], "tables": []}
            
            # 텍스트 추출 (도형 내 텍스트 포함)
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_data["text"].append(shape.text.strip())
                
                # 표 데이터 추출
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    slide_data["tables"].append(table_data)
            
            # 슬라이드 노트 추출 (추가 맥락 파악용)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                slide_data["notes"] = notes.strip()
                
            self.raw_content.append(slide_data)
        
        return self.raw_content

    def generate_scenario_md(self):
        """추출된 원시 데이터를 Gemini에게 전달하여 구조화된 MD 생성"""
        print(f"🔍 PPTX 데이터 분석 중: {self.pptx_path}")
        raw_data = self.extract_text_and_tables()
        
        # 에이전트에게 전달할 프롬프트
        prompt = f"""
        당신은 'QA 시나리오 분석 전문가'입니다. 
        다음은 웹 서비스 기획서(PPTX)에서 추출된 원시 데이터입니다. 
        이 내용을 바탕으로 테스트 에이전트가 이해하기 쉬운 '비즈니스 시나리오 Markdown'을 작성하세요.

        [추출 데이터]
        {raw_data}

        [작성 규칙]
        1. 전체 기능을 'Feature' 단위로 묶으세요.
        2. 각 단계는 '사용자 액션'과 '기대 결과(Expected Result)'가 명확히 드러나야 합니다.
        3. 표 데이터가 있다면 마크다운 표로 변환하여 포함하세요.
        4. 출력은 순수 Markdown 형식으로만 작성하세요.
        """

        response = model.generate_content(prompt)
        return response.text

# --- 실행부 ---
if __name__ == "__main__":
    # 실행 전 inputs 폴더에 테스트용 pptx가 있어야 합니다.
    input_file = "inputs/test_scenario.pptx"
    output_path = "outputs/scenario_draft.md"

    if os.path.exists(input_file):
        analyst = DocAnalyst(input_file)
        markdown_result = analyst.generate_scenario_md()
        
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(markdown_result)
        
        print(f"✨ 시나리오 생성 완료: {output_path}")
    else:
        print(f"❌ 파일을 찾을 수 없습니다: {input_file}")